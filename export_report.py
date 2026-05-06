"""Export reporting data as PDF or PPTX for leadership sharing."""

import io
import os
from datetime import datetime

# ── PDF (ReportLab) ──────────────────────────────────────────────────────────

def generate_pdf(data, sections):
    """Generate a professional PDF report with KPI cards and data tables."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import mm, cm
    from reportlab.lib.colors import HexColor, white, black
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()

    # Custom styles
    PRIMARY = HexColor("#5046e5")
    DARK = HexColor("#1e293b")
    MUTED = HexColor("#64748b")
    SUCCESS = HexColor("#10b981")
    WARNING = HexColor("#f59e0b")
    DANGER = HexColor("#ef4444")
    LIGHT_BG = HexColor("#f8f9fb")

    title_style = ParagraphStyle('ReportTitle', parent=styles['Title'],
                                  fontSize=22, textColor=DARK, spaceAfter=4,
                                  fontName='Helvetica-Bold')
    subtitle_style = ParagraphStyle('ReportSubtitle', parent=styles['Normal'],
                                     fontSize=10, textColor=MUTED, spaceAfter=20)
    heading_style = ParagraphStyle('SectionHeading', parent=styles['Heading2'],
                                    fontSize=13, textColor=PRIMARY, spaceBefore=18, spaceAfter=10,
                                    fontName='Helvetica-Bold')
    metric_label_style = ParagraphStyle('MetricLabel', parent=styles['Normal'],
                                         fontSize=7, textColor=MUTED, fontName='Helvetica')
    metric_value_style = ParagraphStyle('MetricValue', parent=styles['Normal'],
                                         fontSize=16, textColor=DARK, fontName='Helvetica-Bold')
    body_style = ParagraphStyle('Body', parent=styles['Normal'],
                                 fontSize=9, textColor=DARK, leading=13)
    small_style = ParagraphStyle('Small', parent=styles['Normal'],
                                  fontSize=8, textColor=MUTED)

    story = []

    # ── Header ──
    story.append(Paragraph("BSO LUX — Support Analytics Report", title_style))
    date_str = datetime.now().strftime("%B %d, %Y")
    filters_desc = ""
    if data.get("filter_client"):
        filters_desc += f" | Client: {data['filter_client']}"
    if data.get("filter_date_from") or data.get("filter_date_to"):
        filters_desc += f" | Period: {data.get('filter_date_from', '—')} to {data.get('filter_date_to', '—')}"
    story.append(Paragraph(f"Generated: {date_str}{filters_desc}", subtitle_style))

    # ── Key Metrics ──
    if "key_metrics" in sections:
        story.append(Paragraph("Key Performance Indicators", heading_style))

        metrics = [
            ("Total Tickets", str(data["total"]), None),
            ("Open / In Progress", str(data["open_count"]), WARNING),
            ("Resolved / Closed", str(data["resolved_count"]), SUCCESS),
            ("Avg. Resolution", f"{data['avg_resolution_days']}d", None),
            ("Avg. RICE Score", str(data["avg_rice"]), None),
            ("Avg. First Response", f"{data['avg_first_response_hours']}h", None),
            ("SLA Resolution", f"{data['sla_resolution_pct']}%", SUCCESS if data['sla_resolution_pct'] >= 80 else DANGER),
            ("First Response SLA", f"{data['sla_response_pct']}%", SUCCESS if data['sla_response_pct'] >= 80 else DANGER),
            ("This Week", str(data["tickets_this_week"]), None),
        ]

        # Build a 3-column metric card table
        rows = []
        row = []
        for i, (label, value, color) in enumerate(metrics):
            cell_content = [
                Paragraph(label, metric_label_style),
                Paragraph(value, ParagraphStyle('mv', parent=metric_value_style,
                                                 textColor=color or PRIMARY))
            ]
            row.append(cell_content)
            if len(row) == 3 or i == len(metrics) - 1:
                while len(row) < 3:
                    row.append("")
                rows.append(row)
                row = []

        if rows:
            col_width = (A4[0] - 40*mm) / 3
            t = Table(rows, colWidths=[col_width]*3)
            t.setStyle(TableStyle([
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                ('TOPPADDING', (0, 0), (-1, -1), 8),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
                ('LEFTPADDING', (0, 0), (-1, -1), 10),
                ('RIGHTPADDING', (0, 0), (-1, -1), 10),
                ('BACKGROUND', (0, 0), (-1, -1), LIGHT_BG),
                ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
                ('INNERGRID', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
            ]))
            story.append(t)
            story.append(Spacer(1, 10))

    # ── Distribution Tables ──
    def add_dist_table(title, dist_data, key_name, section_id):
        if section_id not in sections or not dist_data:
            return
        story.append(Paragraph(title, heading_style))
        table_data = [["Category", "Count", "Share"]]
        total = sum(d.get("count", d.get(1, 0)) if isinstance(d, dict) else d[1] for d in dist_data)
        for d in dist_data:
            if isinstance(d, dict):
                name = d.get(key_name, d.get("status", "Unknown"))
                count = d.get("count", 0)
            else:
                name = d[0]
                count = d[1]
            pct = f"{(count/total*100):.1f}%" if total > 0 else "0%"
            table_data.append([str(name), str(count), pct])

        t = Table(table_data, colWidths=[90*mm, 35*mm, 35*mm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), PRIMARY),
            ('TEXTCOLOR', (0, 0), (-1, 0), white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('LEFTPADDING', (0, 0), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
            ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
            ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    add_dist_table("Ticket Status Distribution", data.get("status_dist", []), "status", "status_chart")
    add_dist_table("Classification Breakdown", data.get("class_dist", []), "classification", "classification_chart")
    add_dist_table("Risk Level Distribution", data.get("risk_dist", []), "risk_level", "risk_chart")
    add_dist_table("PO Decisions", data.get("po_dist", []), "po_decision", "po_decisions")

    # ── Resolution Time Buckets ──
    if "resolution_time" in sections and data.get("res_buckets"):
        story.append(Paragraph("Time to Resolution", heading_style))
        table_data = [["Bucket", "Count"]]
        for bucket, count in data["res_buckets"].items():
            if count > 0:
                table_data.append([bucket, str(count)])
        if len(table_data) > 1:
            t = Table(table_data, colWidths=[100*mm, 60*mm])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), PRIMARY),
                ('TEXTCOLOR', (0, 0), (-1, 0), white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
                ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
                ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
            ]))
            story.append(t)
            story.append(Spacer(1, 8))

    # ── First Response Time ──
    if "first_response" in sections and data.get("first_response_buckets"):
        story.append(Paragraph("First Response Time Distribution", heading_style))
        table_data = [["Bucket", "Count"]]
        for bucket, count in data["first_response_buckets"].items():
            if count > 0:
                table_data.append([bucket, str(count)])
        if len(table_data) > 1:
            t = Table(table_data, colWidths=[100*mm, 60*mm])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), HexColor("#0891b2")),
                ('TEXTCOLOR', (0, 0), (-1, 0), white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
                ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
                ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
            ]))
            story.append(t)
            story.append(Spacer(1, 8))

    # ── RICE Score Buckets ──
    if "rice_chart" in sections and data.get("rice_buckets"):
        story.append(Paragraph("RICE Score Distribution", heading_style))
        table_data = [["Score Range", "Count"]]
        for bucket, count in data["rice_buckets"].items():
            if count > 0:
                table_data.append([bucket, str(count)])
        if len(table_data) > 1:
            t = Table(table_data, colWidths=[100*mm, 60*mm])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), HexColor("#7c3aed")),
                ('TEXTCOLOR', (0, 0), (-1, 0), white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
                ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
                ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
            ]))
            story.append(t)
            story.append(Spacer(1, 8))

    # ── Templates ──
    if "templates" in sections and data.get("template_breakdown"):
        story.append(Paragraph("Most Reported Templates", heading_style))
        table_data = [["Template", "Tickets"]]
        for item in data["template_breakdown"]:
            table_data.append([item[0], str(item[1])])
        t = Table(table_data, colWidths=[120*mm, 40*mm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), HexColor("#059669")),
            ('TEXTCOLOR', (0, 0), (-1, 0), white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
            ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
            ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    # ── Workflows ──
    if "workflows" in sections and data.get("workflow_breakdown"):
        story.append(Paragraph("Tickets by Workflow", heading_style))
        table_data = [["Workflow", "Tickets"]]
        for item in data["workflow_breakdown"]:
            table_data.append([item[0], str(item[1])])
        t = Table(table_data, colWidths=[120*mm, 40*mm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), HexColor("#0369a1")),
            ('TEXTCOLOR', (0, 0), (-1, 0), white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
            ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
            ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    # ── Timeline ──
    if "timeline" in sections and data.get("monthly"):
        story.append(Paragraph("Tickets Created Over Time (Monthly)", heading_style))
        table_data = [["Month", "Tickets"]]
        for m in data["monthly"]:
            table_data.append([m["month"], str(m["count"])])
        if len(table_data) > 1:
            t = Table(table_data, colWidths=[100*mm, 60*mm])
            t.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), PRIMARY),
                ('TEXTCOLOR', (0, 0), (-1, 0), white),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
                ('TOPPADDING', (0, 0), (-1, -1), 6),
                ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
                ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
                ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
            ]))
            story.append(t)
            story.append(Spacer(1, 8))

    # ── Client Breakdown ──
    if "client_breakdown" in sections and data.get("client_breakdown"):
        story.append(Paragraph("Client / Company Breakdown", heading_style))
        table_data = [["#", "Client", "Tickets", "Share"]]
        total = data["total"] or 1
        for i, (name, count) in enumerate(data["client_breakdown"]):
            pct = f"{(count/total*100):.1f}%"
            table_data.append([str(i+1), name, str(count), pct])
        t = Table(table_data, colWidths=[15*mm, 85*mm, 30*mm, 30*mm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), DARK),
            ('TEXTCOLOR', (0, 0), (-1, 0), white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('ALIGN', (0, 0), (0, -1), 'CENTER'),
            ('ALIGN', (2, 0), (-1, -1), 'CENTER'),
            ('TOPPADDING', (0, 0), (-1, -1), 6),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
            ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
            ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
        ]))
        story.append(t)
        story.append(Spacer(1, 8))

    # ── SLA Compliance ──
    if "sla_compliance" in sections:
        story.append(Paragraph("SLA Compliance Summary", heading_style))
        sla_data = [
            ["Metric", "Target", "Actual", "Status"],
            ["Resolution within 7 days", "80%", f"{data['sla_resolution_pct']}%",
             "PASS" if data['sla_resolution_pct'] >= 80 else "FAIL"],
            ["First response within 24h", "80%", f"{data['sla_response_pct']}%",
             "PASS" if data['sla_response_pct'] >= 80 else "FAIL"],
        ]
        t = Table(sla_data, colWidths=[55*mm, 30*mm, 35*mm, 40*mm])
        t.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), DARK),
            ('TEXTCOLOR', (0, 0), (-1, 0), white),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, -1), 9),
            ('ALIGN', (1, 0), (-1, -1), 'CENTER'),
            ('TOPPADDING', (0, 0), (-1, -1), 8),
            ('BOTTOMPADDING', (0, 0), (-1, -1), 8),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [white, LIGHT_BG]),
            ('BOX', (0, 0), (-1, -1), 0.5, HexColor("#e2e8f0")),
            ('INNERGRID', (0, 0), (-1, -1), 0.3, HexColor("#e2e8f0")),
        ]))
        story.append(t)

    # ── Footer ──
    story.append(Spacer(1, 20))
    story.append(Paragraph("Generated by Freshdesk AI Analyzer — BSO LUX Templates Team", small_style))

    doc.build(story)
    buf.seek(0)
    return buf


# ── PPTX (python-pptx) ──────────────────────────────────────────────────────

def generate_pptx(data, sections):
    """Generate a professional PPTX presentation for leadership."""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.chart import XL_CHART_TYPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Colors
    PRIMARY = RGBColor(0x50, 0x46, 0xe5)
    DARK = RGBColor(0x1e, 0x29, 0x3b)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MUTED = RGBColor(0x64, 0x74, 0x8b)
    SUCCESS = RGBColor(0x10, 0xb9, 0x81)
    WARNING = RGBColor(0xf5, 0x9e, 0x0b)
    ACCENT = RGBColor(0x06, 0xb6, 0xd4)
    LIGHT_BG = RGBColor(0xf8, 0xf9, 0xfb)

    SLIDE_W = Inches(13.333)
    SLIDE_H = Inches(7.5)

    def add_bg(slide, color):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_text_box(slide, left, top, width, height, text, font_size=12,
                     bold=False, color=DARK, align=PP_ALIGN.LEFT, font_name="Calibri"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = font_name
        p.alignment = align
        return txBox

    # ── Slide 1: Title ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide, DARK)

    # Accent bar
    shape = slide.shapes.add_shape(1, Inches(0), Inches(3.2), Inches(0.15), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    add_text_box(slide, Inches(0.8), Inches(2.2), Inches(10), Inches(1),
                 "BSO LUX — Support Analytics", 38, True, WHITE, font_name="Calibri")

    date_str = datetime.now().strftime("%B %d, %Y")
    filters_text = f"Report generated: {date_str}"
    if data.get("filter_client"):
        filters_text += f"  |  Client: {data['filter_client']}"
    if data.get("filter_date_from") or data.get("filter_date_to"):
        filters_text += f"  |  Period: {data.get('filter_date_from', '—')} to {data.get('filter_date_to', '—')}"
    add_text_box(slide, Inches(0.8), Inches(3.5), Inches(10), Inches(0.5),
                 filters_text, 14, False, MUTED, font_name="Calibri")

    add_text_box(slide, Inches(0.8), Inches(5.5), Inches(10), Inches(0.4),
                 "Silverfin Luxembourg Templates Team", 12, False, RGBColor(0x94, 0xa3, 0xb8), font_name="Calibri")

    # ── Slide 2: Key Metrics ──
    if "key_metrics" in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)

        add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8), Inches(0.6),
                     "Key Performance Indicators", 24, True, DARK, font_name="Calibri")

        # Accent line under title
        shape = slide.shapes.add_shape(1, Inches(0.6), Inches(1.0), Inches(2), Inches(0.04))
        shape.fill.solid()
        shape.fill.fore_color.rgb = PRIMARY
        shape.line.fill.background()

        metrics = [
            ("Total Tickets", str(data["total"]), PRIMARY),
            ("Open / In Progress", str(data["open_count"]), WARNING),
            ("Resolved / Closed", str(data["resolved_count"]), SUCCESS),
            ("Avg. Resolution", f"{data['avg_resolution_days']} days", ACCENT),
            ("Avg. RICE Score", str(data["avg_rice"]), RGBColor(0x7c, 0x3a, 0xed)),
            ("Avg. First Response", f"{data['avg_first_response_hours']}h", ACCENT),
            ("SLA Resolution", f"{data['sla_resolution_pct']}%", SUCCESS if data['sla_resolution_pct'] >= 80 else RGBColor(0xef, 0x44, 0x44)),
            ("First Response SLA", f"{data['sla_response_pct']}%", SUCCESS if data['sla_response_pct'] >= 80 else RGBColor(0xef, 0x44, 0x44)),
            ("This Week", str(data["tickets_this_week"]), RGBColor(0x8b, 0x5c, 0xf6)),
        ]

        cols = 3
        card_w = Inches(3.8)
        card_h = Inches(1.4)
        gap_x = Inches(0.35)
        gap_y = Inches(0.2)
        start_x = Inches(0.6)
        start_y = Inches(1.4)

        for i, (label, value, color) in enumerate(metrics):
            col = i % cols
            row = i // cols
            x = start_x + col * (card_w + gap_x)
            y = start_y + row * (card_h + gap_y)

            # Card background
            card = slide.shapes.add_shape(1, x, y, card_w, card_h)
            card.fill.solid()
            card.fill.fore_color.rgb = LIGHT_BG
            card.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)
            card.line.width = Pt(1)

            # Left accent
            accent = slide.shapes.add_shape(1, x, y, Inches(0.06), card_h)
            accent.fill.solid()
            accent.fill.fore_color.rgb = color
            accent.line.fill.background()

            add_text_box(slide, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.4), Inches(0.3),
                         label, 10, False, MUTED, font_name="Calibri")
            add_text_box(slide, x + Inches(0.25), y + Inches(0.55), card_w - Inches(0.4), Inches(0.6),
                         value, 28, True, color, font_name="Calibri")

    # ── Helper: data table slide ──
    def add_table_slide(title, headers, rows, header_color=PRIMARY):
        if not rows:
            return
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)

        add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
                     title, 24, True, DARK, font_name="Calibri")

        n_rows = len(rows) + 1
        n_cols = len(headers)
        table_width = min(Inches(12), SLIDE_W - Inches(1.2))
        col_w = table_width / n_cols

        tbl = slide.shapes.add_table(n_rows, n_cols,
                                      Inches(0.6), Inches(1.3),
                                      table_width, Inches(min(5.5, 0.45 * n_rows))).table

        # Header
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_color
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            if j > 0:
                p.alignment = PP_ALIGN.CENTER

        # Data
        for i, row in enumerate(rows):
            for j, val in enumerate(row):
                cell = tbl.cell(i + 1, j)
                cell.text = str(val)
                if i % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = LIGHT_BG
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(10)
                p.font.color.rgb = DARK
                p.font.name = "Calibri"
                if j > 0:
                    p.alignment = PP_ALIGN.CENTER

    # ── Distribution slides ──
    def dist_to_rows(dist_data, key_name):
        total = sum(d.get("count", 0) for d in dist_data) if dist_data else 0
        rows = []
        for d in dist_data:
            name = d.get(key_name, "Unknown")
            count = d.get("count", 0)
            pct = f"{(count/total*100):.1f}%" if total > 0 else "0%"
            rows.append([name, count, pct])
        return rows

    if "status_chart" in sections and data.get("status_dist"):
        add_table_slide("Ticket Status Distribution",
                        ["Status", "Count", "Share"],
                        dist_to_rows(data["status_dist"], "status"))

    if "classification_chart" in sections and data.get("class_dist"):
        add_table_slide("Classification Breakdown",
                        ["Classification", "Count", "Share"],
                        dist_to_rows(data["class_dist"], "classification"))

    if "risk_chart" in sections and data.get("risk_dist"):
        add_table_slide("Risk Level Distribution",
                        ["Risk Level", "Count", "Share"],
                        dist_to_rows(data["risk_dist"], "risk_level"),
                        RGBColor(0xef, 0x44, 0x44))

    if "po_decisions" in sections and data.get("po_dist"):
        add_table_slide("PO Decisions",
                        ["Decision", "Count", "Share"],
                        dist_to_rows(data["po_dist"], "po_decision"))

    # ── Resolution & Response buckets ──
    if "resolution_time" in sections and data.get("res_buckets"):
        rows = [[b, c] for b, c in data["res_buckets"].items() if c > 0]
        add_table_slide("Time to Resolution", ["Bucket", "Count"], rows, ACCENT)

    if "first_response" in sections and data.get("first_response_buckets"):
        rows = [[b, c] for b, c in data["first_response_buckets"].items() if c > 0]
        add_table_slide("First Response Time", ["Bucket", "Count"], rows,
                        RGBColor(0x08, 0x91, 0xb2))

    if "rice_chart" in sections and data.get("rice_buckets"):
        rows = [[b, c] for b, c in data["rice_buckets"].items() if c > 0]
        add_table_slide("RICE Score Distribution", ["Range", "Count"], rows,
                        RGBColor(0x7c, 0x3a, 0xed))

    # ── Templates & Workflows ──
    if "templates" in sections and data.get("template_breakdown"):
        rows = [[t[0], t[1]] for t in data["template_breakdown"]]
        add_table_slide("Most Reported Templates", ["Template", "Tickets"], rows,
                        RGBColor(0x05, 0x96, 0x69))

    if "workflows" in sections and data.get("workflow_breakdown"):
        rows = [[w[0], w[1]] for w in data["workflow_breakdown"]]
        add_table_slide("Tickets by Workflow", ["Workflow", "Tickets"], rows,
                        RGBColor(0x03, 0x69, 0xa1))

    # ── Timeline ──
    if "timeline" in sections and data.get("monthly"):
        rows = [[m["month"], m["count"]] for m in data["monthly"]]
        add_table_slide("Tickets Created Over Time", ["Month", "Tickets"], rows)

    # ── Client Breakdown ──
    if "client_breakdown" in sections and data.get("client_breakdown"):
        total = data["total"] or 1
        rows = [[i+1, name, count, f"{(count/total*100):.1f}%"]
                for i, (name, count) in enumerate(data["client_breakdown"])]
        add_table_slide("Client Breakdown", ["#", "Client", "Tickets", "Share"], rows, DARK)

    # ── SLA Compliance ──
    if "sla_compliance" in sections:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_bg(slide, WHITE)

        add_text_box(slide, Inches(0.6), Inches(0.4), Inches(10), Inches(0.6),
                     "SLA Compliance", 24, True, DARK, font_name="Calibri")

        # Resolution SLA card
        res_pass = data['sla_resolution_pct'] >= 80
        card = slide.shapes.add_shape(1, Inches(0.6), Inches(1.5), Inches(5.5), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)

        add_text_box(slide, Inches(1), Inches(1.7), Inches(4.5), Inches(0.4),
                     "Resolution within 7 days", 14, True, DARK, font_name="Calibri")
        add_text_box(slide, Inches(1), Inches(2.2), Inches(4.5), Inches(1),
                     f"{data['sla_resolution_pct']}%", 48, True,
                     SUCCESS if res_pass else RGBColor(0xef, 0x44, 0x44), font_name="Calibri")
        add_text_box(slide, Inches(1), Inches(3.3), Inches(4.5), Inches(0.4),
                     "PASS" if res_pass else "NEEDS IMPROVEMENT", 12, True,
                     SUCCESS if res_pass else RGBColor(0xef, 0x44, 0x44), font_name="Calibri")

        # First Response SLA card
        resp_pass = data['sla_response_pct'] >= 80
        card2 = slide.shapes.add_shape(1, Inches(6.8), Inches(1.5), Inches(5.5), Inches(2.5))
        card2.fill.solid()
        card2.fill.fore_color.rgb = LIGHT_BG
        card2.line.color.rgb = RGBColor(0xe2, 0xe8, 0xf0)

        add_text_box(slide, Inches(7.2), Inches(1.7), Inches(4.5), Inches(0.4),
                     "First Response within 24h", 14, True, DARK, font_name="Calibri")
        add_text_box(slide, Inches(7.2), Inches(2.2), Inches(4.5), Inches(1),
                     f"{data['sla_response_pct']}%", 48, True,
                     SUCCESS if resp_pass else RGBColor(0xef, 0x44, 0x44), font_name="Calibri")
        add_text_box(slide, Inches(7.2), Inches(3.3), Inches(4.5), Inches(0.4),
                     "PASS" if resp_pass else "NEEDS IMPROVEMENT", 12, True,
                     SUCCESS if resp_pass else RGBColor(0xef, 0x44, 0x44), font_name="Calibri")

    # ── Final slide ──
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK)
    add_text_box(slide, Inches(0.8), Inches(3), Inches(10), Inches(0.8),
                 "Thank you", 36, True, WHITE, PP_ALIGN.LEFT, "Calibri")
    add_text_box(slide, Inches(0.8), Inches(3.8), Inches(10), Inches(0.5),
                 "BSO LUX Templates Team — Freshdesk AI Analyzer", 14, False, MUTED,
                 PP_ALIGN.LEFT, "Calibri")

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
